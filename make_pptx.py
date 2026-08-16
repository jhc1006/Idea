import sys
import os
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    NAVY = RGBColor(6, 9, 19)
    WHITE = RGBColor(248, 250, 252)
    GREY = RGBColor(148, 163, 184)
    CYAN = RGBColor(56, 189, 248)
    GREEN = RGBColor(16, 185, 129)

    blank_layout = prs.slide_layouts[6]

    # Helper function to set solid background
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1, NAVY)

    tb = slide1.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "생성형 AI 및 IoT 텔레메트리 기반 독거노인 관제 & 사회복지사 실시간 연계 플랫폼"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Outfit"

    p2 = tf.add_paragraph()
    p2.text = "SilverCare AI"
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "Outfit"
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "초고령사회 고독사 예방을 위한 스마트 골든타임 케어 솔루션"
    p3.font.size = Pt(20)
    p3.font.color.rgb = GREY
    p3.font.name = "Inter"
    p3.space_before = Pt(10)

    p4 = tf.add_paragraph()
    p4.text = "발표자: 전효철 (AI & Full-Stack Software Engineer)\n시연 주소: http://localhost:8085"
    p4.font.size = Pt(14)
    p4.font.color.rgb = CYAN
    p4.font.name = "Inter"
    p4.space_before = Pt(30)

    slide1.notes_slide.notes_text_frame.text = (
        "안녕하십니까, 'SilverCare AI' 사업 과제 발표를 맡은 주관기관 대표 전효철입니다. "
        "오늘 저희가 선보일 솔루션은 AI 텔레메트리 기술을 활용하여 초고령사회 독거노인의 고독사를 예방하고, "
        "위급 상황 발생 시 3분 이내에 사회복지사가 현장에 출동하도록 연계하는 스마트 관제 플랫폼입니다."
    )

    # Slide 2: Problem
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2, NAVY)

    tb = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "추진 배경 및 해결하고자 하는 문제 (Problem)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"

    problems = [
      ("사회복지사 인력 부족", "복지사 1인당 30~50가구 담당\n업무 과중으로 밀착 모니터링 및 실시간 돌봄 공백 발생"),
      ("단순 센서의 높은 오작동", "기존 센서의 빈번한 오탐지로 관제 피로 누적 및 쓰러짐 등 실제 긴급상황 골든타임 감지 지연"),
      ("데이터의 연동성 부재", "관제 센터 경보 발생 시에도 현장 복지사 및 보호자와의 데이터 단절로 출동 및 확인 조치 지연")
    ]

    for i, (title, desc) in enumerate(problems):
        left = Inches(1 + i * 3.8)
        top = Inches(1.8)
        width = Inches(3.5)
        height = Inches(4.5)

        shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(13, 21, 39)
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)

        tb_card = slide2.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf_card = tb_card.text_frame
        tf_card.word_wrap = True

        p_title = tf_card.paragraphs[0]
        p_title.text = f"0{i+1}. {title}"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = CYAN
        p_title.font.name = "Outfit"

        p_desc = tf_card.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = GREY
        p_desc.font.name = "Inter"
        p_desc.space_before = Pt(20)

    slide2.notes_slide.notes_text_frame.text = (
        "2026년 대한민국은 65세 이상 인구가 20%를 넘어서는 초고령사회에 접어들었습니다. 독거노인 가구는 200만을 넘었지만, "
        "현장의 사회복지사 한 명이 담당해야 하는 어르신은 40가구에 달합니다. "
        "기존 센서는 오작동이 많고 쓰러짐이나 질환 발생 시 관제 센터와 현장 출동 복지사 간에 데이터가 즉시 공유되지 않아 "
        "골든타임을 놓치는 비극이 반복되고 있습니다."
    )

    # Slide 3: Solution
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3, NAVY)

    tb = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "SilverCare AI 4단계 원스톱 케어 파이프라인 (Solution)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"

    steps = [
      ("1단계: 실시간 감지", "댁내 움직임 + 스마트 약통 + AI 스피커 대화 맥락 + SOS 버튼 텔레메트리 통합 수집"),
      ("2단계: AI 위험 분석", "AI 분석기가 위험도 4단계(정상, 주의, 경고, 긴급)로 자동 검증 및 분석 수행"),
      ("3단계: 센터 관제 팝업", "서울시 마포구 관제 센터 모니터에 시각 경보 및 위험 인물 자동 하이라이트"),
      ("4단계: 사회복지사 출동", "담당 복지사 모바일 앱에 기저질환과 현장 위치 자동 매칭 전송 및 보호자 카톡 알림")
    ]

    for i, (title, desc) in enumerate(steps):
        left = Inches(1 + i * 2.85)
        top = Inches(2.2)
        width = Inches(2.7)
        height = Inches(3.8)

        shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(13, 21, 39)
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)

        tb_step = slide3.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True

        p_title = tf_step.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = CYAN
        p_title.font.name = "Outfit"

        p_desc = tf_step.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = GREY
        p_desc.font.name = "Inter"
        p_desc.space_before = Pt(20)

    slide3.notes_slide.notes_text_frame.text = (
        "저희 SilverCare AI는 이 문제를 해결하기 위해 4단계 통합 파이프라인을 구축했습니다. "
        "어르신 댁내 멀티 센서 데이터를 수집하여 AI가 위험도를 4단계로 자동 판단하고, "
        "관제 센터 팝업과 동시에 담당 사회복지사 스마트폰으로 위험 정보와 어르신 기저질환을 즉시 전송합니다. "
        "현장 출동부터 보호자 통보까지 단 3분 만에 완료됩니다."
    )

    # Slide 4: Core Tech & Live Demo
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4, NAVY)

    tb = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "핵심 기술 경쟁력 및 프로토타입 시연"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"

    tb_details = slide4.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(4.5))
    tf_details = tb_details.text_frame
    tf_details.word_wrap = True

    bullets = [
      ("IoT 멀티 텔레메트리 융합 분석 기술", "24시간 활동 무반응, 연속 복약 미이행, AI 대화 중 통증/외로움 키워드 정밀 감지"),
      ("현장 사회복지사 연계 모바일 워크플로우", "긴급상황 발생 시 담당 구역 복지사에 푸시 알림, 출동 지시 및 조치 보고 시 즉시 관제판 해제 연동"),
      ("SLA 벤치마크 500ms 만족", "실시간 비동기 패킷 파이프라인 설계를 통해 대기 시간(Latency) 최소화 구현 완료"),
      ("프로토타입 실시 시연 (http://localhost:8085)", "로컬 서버 웹 대시보드와 백엔드 모의 텔레메트리 엔진(mock_telemetry.py) 연동 상태 시연 가능")
    ]

    for title, desc in bullets:
        p_t = tf_details.add_paragraph()
        p_t.text = f"⚡ {title}"
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.font.name = "Outfit"
        p_t.space_before = Pt(10)

        p_d = tf_details.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = GREY
        p_d.font.name = "Inter"
        p_d.space_before = Pt(4)
        p_d.level = 0

    slide4.notes_slide.notes_text_frame.text = (
        "화면에 보이시는 대시보드가 저희가 실제 구현한 SilverCare 관제 센터 프로토타입입니다. "
        "우측 상단의 'AI 센서 모의 발생기'를 클릭하면 어르신의 SOS 버튼 작동이나 24시간 활동 중단 신호가 실시간 텔레메트리로 전송되고, "
        "관제판에 긴급 경보와 함께 담당 복지사의 출동 지시 버튼이 활성화됩니다. 조치 완료 시 보호자에게도 자동 알림이 전송됩니다."
    )

    # Slide 5: Business Model & Growth Roadmap
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5, NAVY)

    tb = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "비즈니스 모델 및 성장 로드맵"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"

    bms = [
      ("B2G (공공 돌봄)", "지자체/보건소/복지관 대상 구축비 + 가구당 월 관제 리스료 (12,000원/월)"),
      ("B2C (민간 돌봄)", "타지 거주 자녀 대상 프리미엄 안심 케어 구독 (19,900원/월, 주간 보고서 제공)"),
      ("B2B (실버타운 연계)", "요양병원, 실버타운, 요양원 내 API 및 데이터 연동 라이선스 모델")
    ]

    for i, (title, desc) in enumerate(bms):
        left = Inches(1)
        top = Inches(1.8 + i * 1.6)
        width = Inches(11.3)
        height = Inches(1.3)

        shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(13, 21, 39)
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)

        tb_bm = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
        tf_bm = tb_bm.text_frame
        tf_bm.word_wrap = True

        p_t = tf_bm.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.font.name = "Outfit"

        p_d = tf_bm.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = GREY
        p_d.font.name = "Inter"
        p_d.space_before = Pt(6)

    slide5.notes_slide.notes_text_frame.text = (
        "저희 비즈니스 모델은 지자체 고독사 예방 예산을 활용하는 B2G 리스 모델과 타지 거주 자녀를 타겟으로 하는 B2C 월 구독 모델로 구성됩니다. "
        "1년차 마포구 150가구 시범 사업을 시작으로 3년 차에는 전국 229개 시군구 확장을 통해 50억 원의 매출을 달성하겠습니다."
    )

    # Slide 6: Team & Budget
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6, NAVY)

    tb = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "팀 역량 및 2억원 사업비 집행 계획 (Team & Budget)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Outfit"

    budgets = [
      ("인건비 (40%)", "80,000,000원", "AI 관제 알고리즘, 실시간 웹서버 및 사회복지사 연계 모바일 연동 인력 구성"),
      ("연구장비 및 재료비 (35%)", "70,000,000원", "어르신 댁내 IoT 센서(약통, 움직임 감지기), 모의 하네스 구축 장비비"),
      ("사업화 및 마케팅비 (15%)", "30,000,000원", "서울시 마포구 PoC 시범 사업 및 지자체 스마트복지 바우처 판로 개척비"),
      ("위탁개발 및 자문비 (10%)", "20,000,000원", "사회복지 전문가 자문, 데이터 프라이버시 및 보안 법률 가이드 확보비")
    ]

    for i, (title, val, desc) in enumerate(budgets):
        left = Inches(1 + i * 2.85)
        top = Inches(1.8)
        width = Inches(2.7)
        height = Inches(4.5)

        shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(13, 21, 39)
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)

        tb_b = slide6.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True

        p_t = tf_b.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.font.name = "Outfit"

        p_v = tf_b.add_paragraph()
        p_v.text = val
        p_v.font.size = Pt(18)
        p_v.font.bold = True
        p_v.font.color.rgb = GREEN
        p_v.font.name = "Outfit"
        p_v.space_before = Pt(10)

        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = GREY
        p_d.font.name = "Inter"
        p_d.space_before = Pt(15)

    slide6.notes_slide.notes_text_frame.text = (
        "저는 생성형 AI와 멀티에이전트, 고성능 대시보드 시스템을 직접 설계하고 구현할 수 있는 풀스택 개발 역량을 갖추고 있습니다. "
        "총 2억 원의 사업비 중 75%를 핵심 개발 인력과 센서 인프라에 투입하여 12개월 이내에 완성도 높은 시스템을 납품하겠습니다."
    )

    # Slide 7: Closing
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7, NAVY)

    tb = slide7.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "기술로 단 한 명의 어르신도 방치되지 않는 따뜻한 사회"
    p.font.size = Pt(28)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.font.name = "Outfit"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "경청해 주셔서 감사합니다. 질의응답에 성심껏 답변드리겠습니다."
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.font.name = "Inter"
    p2.space_before = Pt(30)
    p2.alignment = PP_ALIGN.CENTER

    slide7.notes_slide.notes_text_frame.text = (
        "SilverCare AI는 단순한 기술 개발을 넘어 단 한 명의 어르신도 홀로 방치되지 않는 따뜻한 사회적 안전망을 만드는 과제입니다. "
        "경청해 주셔서 감사합니다. 질의응답에 성심껏 답변드리겠습니다."
    )

    prs.save("/Users/jeonhyochul/work/Idea/SilverCare_AI_Government_Grant_Presentation.pptx")
    print("✨ PPTX pitch deck generated successfully.")

if __name__ == '__main__':
    create_presentation()
